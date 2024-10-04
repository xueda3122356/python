# pip install python-pptx
def base_ppt():
    from pptx import Presentation
    # create power point
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide = ppt.slides.add_slide(ppt.slide_layouts[2])
    slide = ppt.slides.add_slide(ppt.slide_layouts[3])
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])


    ppt.save('./generate_data/26_create_ppt.pptx')


def base_ppt1():
    from pptx import Presentation
    from pptx.util import Pt
    # create power point
    ppt = Presentation()
    # create new slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])

    # select text frame
    shapes = slide.shapes
    title = shapes.title
    title.text = 'This is a practise for learning creating ppt by python'
    content = shapes.placeholders[1]
    content.text = 'Content Info'

    # create list
    slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
    shapes = slide2.shapes
    content = shapes.placeholders[1]
    tf = content.text_frame
    p1 = tf.add_paragraph()
    p2 = tf.add_paragraph()
    p3 = tf.add_paragraph()
    p1.text = 'This is the first paragraph'
    p1.level = 1
    p2.text = 'This is the second paragraph'
    p2.level = 2
    p3.text = 'This is the third paragraph'
    p3.level = 3

    # font setting
    p4 = tf.add_paragraph()
    p4.text = 'This is the fourth paragraph'
    p4.font.bold = True
    p4.font.size = Pt(30)

    # create a new text box
    slide3 = ppt.slides.add_slide(ppt.slide_layouts[6])
    left = top = width = height = Pt(200)
    text_box = slide3.shapes.add_textbox(left,top,width,height)
    tf = text_box.text_frame
    tf.text = 'this is the info of text box'
    p = tf.add_paragraph()
    p.text = 'this is the paragraph'




    ppt.save('./generate_data/26_create_ppt1.pptx')

if __name__ == "__main__":
    base_ppt1()