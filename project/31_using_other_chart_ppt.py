def line_chart():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.util import Inches,Pt
    
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # Encapsulation chart data
    chart_data = CategoryChartData()

    # Grouping data
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']

    # Specify data(adding more groups data)
    chart_data.add_series('series1', (19,26,17,30))
    chart_data.add_series('series2', (23,22,21,25))
    chart_data.add_series('series3', (16,29,24,32))


    # creating chart
    x = y = Inches(1.5)
    height = Inches(4)
    width = Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE,x,y,width,height,chart_data).chart
    chart.chart_style = 10 # 48 different styles for column
    chart.font.size = Pt(10) # set all labels font size

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(20) # set x-axis labels font size
    category_axis.has_major_gridlines = True # if display as table or not

    # set labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.CENTER

    # set legend
    chart.has_legend = True
    chart.legend.font.size = Pt(10)
    chart.legend.position = XL_LEGEND_POSITION.LEFT
    chart.legend.include_in_layout = False

    ppt.save('./generate_data/31_line_chart_style_ppt.pptx')

def pie_chart():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.util import Inches,Pt
    
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # Encapsulation chart data
    chart_data = CategoryChartData()

    # Grouping data
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']

    chart_data.add_series('季度销售比例', (0.27,0.33,0.19,0.21))

    # creating chart
    x = y = Inches(1.5)
    height = Inches(4)
    width = Inches(6)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE,x,y,width,height,chart_data).chart

    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0%'

    ppt.save('./generate_data/31_pie_chart_style_ppt.pptx')

if __name__ =="__main__":
    #line_chart()
    pie_chart()