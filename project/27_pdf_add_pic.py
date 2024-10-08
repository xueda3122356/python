def add_pic():
    from pptx import Presentation
    from pptx.util import Pt
    ppt = Presentation()
    slide1 = ppt.slides.add_slide(ppt.slide_layouts[1])

    # Get tex box
    shapes = slide1.shapes
    left = top = Pt(100)
    shapes.add_picture('./base_data/background.jpg', left=left, top=top)

    ppt.save('./generate_data/27_ppt_add_picture.pptx')

    # Get tex box
    slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
    shapes = slide2.shapes
    left = top = Pt(0)
    height = width = Pt(100)
    shapes.add_picture('./base_data/background.jpg', left=left, top=top, height=height, width=width)

    ppt.save('./generate_data/27_ppt_add_picture.pptx')

if __name__ == "__main__":
    add_pic()