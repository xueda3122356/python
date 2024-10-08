def base_use():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    shapes = slide.shapes
    shapes.title.text = "流程图"


    left = Inches(1)
    top = Inches(3)
    width = Inches(2)
    height = Inches(1)
    t_sh = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PENTAGON, left, top, width, height)
    t_sh.text = "第1步"

    for i in range(2,6):
        left = left + width - Inches(0.4)
        top = Inches(3)
        width = Inches(2)
        height = Inches(1)
        temp_sh = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
        #temp_sh.text = f'第{i}步'
        p = temp_sh.text_frame
        p.text = f'第{i}步' 
        p.fit_text(max_size=10, bold=True, italic=True) # 修改文本格式必须先输入文本


    ppt.save('./generate_data/28_flow_chart_ppt.pptx')
    

if __name__ == "__main__":
    base_use()