def add_chart():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # Encapsulation chart data
    chart_data = CategoryChartData()

    # Grouping data
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']

    # Specify data
    chart_data.add_series('series', (19,26,17,30))


    # creating chart
    x = y = Inches(1.5)
    height = Inches(4.5)
    width = Inches(6)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,x,y,width,height,chart_data)

    ppt.save('./generate_data/29_create_chart_ppt.pptx')


def add_chart2():
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # Encapsulation chart data
    chart_data = CategoryChartData()

    # Grouping data
    chart_data.categories = ['第一季度','第二季度','第三季度','第四季度']

    # Specify data(adding more groups data)
    chart_data.add_series('series', (19,26,17,30))
    chart_data.add_series('series', (23,22,21,25))
    chart_data.add_series('series', (16,29,24,32))


    # creating chart
    x = y = Inches(1.5)
    height = Inches(4.5)
    width = Inches(6)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,x,y,width,height,chart_data)

    ppt.save('./generate_data/29_create_chart_ppt2.pptx')


if __name__ == "__main__":
    #add_chart()
    add_chart2()